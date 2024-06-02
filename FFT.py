from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new PowerPoint presentation
ppt = Presentation()

# Add title slide
slide_layout = ppt.slide_layouts[0]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Fragrance Fusion: An E-commerce Experience"
subtitle.text = "Exploring the Online Marketplace for Fragrances\nYour Name\nDate"


# Define a function to add bullet point slides
def add_bullet_slide(presentation, title, bullet_points):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    tf = body_shape.text_frame
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)


# Add the slides with content
slides = [
    ("Introduction", [
        "Overview of Fragrance Fusion",
        "Online platform specializing in fragrances",
        "Variety of products: perfumes, colognes, essential oils, etc.",
        "Importance of E-commerce in the Fragrance Industry",
        "Convenience and accessibility",
        "Wide product range"
    ]),
    ("Objectives", [
        "Main Goals of the Project",
        "Understand the user experience on Fragrance Fusion",
        "Analyze the effectiveness of the website design",
        "Evaluate the purchasing process",
        "Identify areas for improvement"
    ]),
    ("Literature Review", [
        "Existing Research on E-commerce Websites",
        "User experience and interface design",
        "Online shopping behavior in the fragrance market",
        "Key factors for successful e-commerce platforms",
        "Relevant Theories and Models",
        "Technology Acceptance Model (TAM)",
        "E-commerce adoption frameworks"
    ]),
    ("Methodology/Methods", [
        "Research Approach",
        "Qualitative and quantitative methods",
        "Data Collection Techniques",
        "Surveys and questionnaires",
        "User interviews",
        "Website analytics and usability testing",
        "Sample Size and Demographics",
        "Number of participants",
        "Age, gender, and other relevant demographics"
    ]),
    ("Design", [
        "Website Design Features",
        "User interface (UI) and user experience (UX)",
        "Navigation and layout",
        "Product pages and descriptions",
        "Checkout process",
        "Design Principles Applied",
        "Simplicity and clarity",
        "Responsiveness and accessibility"
    ]),
    ("Implementation", [
        "Steps Taken for Implementation",
        "Development of the website",
        "Integration of e-commerce functionalities",
        "Testing and debugging",
        "Launch and promotion strategies",
        "Tools and Technologies Used",
        "E-commerce platforms (e.g., Shopify, WooCommerce)",
        "Web development technologies (e.g., HTML, CSS, JavaScript)"
    ]),
    ("Results and Discussion", [
        "User Feedback and Satisfaction",
        "Survey and interview results",
        "Key findings on user experience",
        "Website Performance Metrics",
        "Traffic and conversion rates",
        "Sales data and revenue",
        "Comparison with Competitors",
        "Strengths and weaknesses"
    ]),
    ("Conclusion", [
        "Summary of Key Points",
        "Achievements and objectives met",
        "Overall performance of Fragrance Fusion",
        "Future Recommendations",
        "Potential improvements",
        "Next steps for growth and development"
    ]),
    ("References", [
        "Citations and Sources",
        "Academic articles and books",
        "Industry reports and market analysis",
        "Websites and online resources"
    ]),
    ("Questions and Answers", [
        "Invitation for Questions",
        "Encourage audience engagement"
    ])
]

for title, content in slides:
    add_bullet_slide(ppt, title, content)

# Save the presentation to a file
ppt_path = "/mnt/data/Fragrance_Fusion_Presentation_v2.pptx"
ppt.save(ppt_path)

ppt_path

# Users/anupjaiswal/Desktop

